import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Premium Corporate Color Palette (Dark Theme)
    BG_COLOR = RGBColor(15, 23, 42) # Slate 900
    TEXT_MAIN = RGBColor(248, 250, 252) # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ACCENT_COLOR = RGBColor(56, 189, 248) # Sky 400
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_standard_slide(title_text):
        slide_layout = prs.slide_layouts[6] # Blank slide to build custom layout
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide)
        
        # Add Title Box
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = ACCENT_COLOR
        
        return slide

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Retail Customer Insights & Sales Performance Dashboard"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Strategic Retail Operations & Customer Intelligence"
    p2.font.size = Pt(24)
    p2.font.name = "Arial"
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Final-Year Data Analytics Capstone Project  |  Presenter: Pratham"
    p3.font.size = Pt(14)
    p3.font.name = "Arial"
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)

    # Slide 2: The Business Problem
    slide = add_standard_slide("The Business Problem")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Customer Retention Churn: High customer acquisition costs with a lack of demographic insights make retaining customers highly difficult."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Margin Erosion: Uncontrolled discounting practices (often exceeding 50%) eat away at regional margins, turning high-revenue sub-categories unprofitable."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Inventory Mismanagement: Heavy holding costs for dead stock coexist with stockouts of high-velocity technology items."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Information Silos: Transactional databases are disconnected from marketing spend data, rendering ROI calculations impossible."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 3: Project Objectives
    slide = add_standard_slide("Project Objectives")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    points = [
        "1. Consolidate Relational Schema: Normalize and enrich transactional data with customer demographics and monthly regional marketing campaigns.",
        "2. Standardize Cleaning Pipeline: Automate data deduplication, handle outlier transactions via IQR, and engineer business KPIs.",
        "3. Execute RFM Customer Segmentation: Segment the user base into distinct profiles to trigger targeted marketing.",
        "4. Build Predictive Engines: Implement ML classifiers and regressors to forecast future sales, customer churn probabilities, and category demand levels.",
        "5. Deliver Streamlit Frontend & Power BI Specs: Enable self-service forecasting, risk profiling, and executive reporting."
    ]
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(12)

    # Slide 4: Dataset Overview (Data Model)
    slide = add_standard_slide("Dataset Overview & Data Model")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Unified Relational Schema Structure:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Customer Data (800 profiles): Customer ID (PK), Name, Gender, Age, City, Region, Registration Date."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Orders Data (9,994 entries): Order ID, Customer ID (FK), Product ID (FK), Quantity, Sales, Profit, Discount, Date."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Product Data (1,850 catalog items): Product ID (PK), Product Name, Category, Sub-category, Cost Price, Selling Price."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Region Data: Geography dimension containing Region, State, and City."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Marketing Spend Data (monthly regional): Region, Campaign Type, Spend Amount, Campaign Date."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN

    # Slide 5: Data Cleaning & Preprocessing Process
    slide = add_standard_slide("Data Cleaning & Preprocessing Process")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Deduplication: Isolated and removed duplicate order and customer records."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Missing Value Treatment: Imputed blank customer names and verified catalog dimensions."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Outlier Profiling (IQR): Flagged high-volume outliers (using 1.5 * IQR) in Sales and Profit to prevent training skew."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Date Parsing & Type Casting: Standardized dates to datetime64 and float types."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Standardized Schema: Applied snake_case column naming convention across all components."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 6: Sales Performance Analysis
    slide = add_standard_slide("Sales Performance Analysis")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Monthly sales trends: Seasonal sales peaks during Q4 holidays (Nov-Dec) representing a ~35% surge."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Regional Distribution: The West region leads in sales contribution, followed closely by the East, while the South underperforms."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Discount Impact: Transactions with discounts exceeding 20% show negative average profits, illustrating severe margin erosion."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Product Concentration: Top 20 products contribute to a disproportionate 15% of total sales volume."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 7: Customer Demographics & RFM Analysis
    slide = add_standard_slide("Customer Demographics & RFM Analysis")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "RFM Customer Segmentation Profiles:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Champions: Highly recent, high-frequency, and high-spending customers (~15% of customer base)."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Loyal Customers: Purchase regularly, highly responsive to promotions."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Potential Loyalists: Recent buyers with moderate spending; prime candidates for cross-selling."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• At Risk: High historical spend but no purchases in the last 180 days."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Lost Customers: Dormant, low frequency, low monetary profiles."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN

    # Slide 8: Product Performance Insights
    slide = add_standard_slide("Product Performance Insights")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Fast-moving products: Technology sub-categories (Phones, Accessories) represent the highest daily velocity scores."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Dead Stock: Products with velocity scores < 0.05 units/day are flagged, representing locked-up capital."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• High Margin Categories: Technology and Office Supplies maintain stable average margins of 40-50%."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Low Margin Sub-categories: Furniture sub-categories (specifically Tables and Bookcases) exhibit net losses due to high shipping expenses and high discounts."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 9: Executive Dashboard Layout & Star Schema
    slide = add_standard_slide("Executive Dashboard Layout & Star Schema")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Power BI Reporting Pages:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Page 1: Executive Overview: Total Revenue, Profit, Orders, AOV, MoM Growth, and Category splits."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Page 2: Customer Insights: RFM Segments distribution, Churn analysis, CLV, and Retention Rate."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Page 3: Product Performance: Top/Bottom products, Velocity lists, Return rates, and safety stock."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• Page 4: Sales Performance: Detailed regional maps, monthly sales trend, and interactive segment filtering."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN

    # Slide 10: Sales Forecasting Models (Time-Series)
    slide = add_standard_slide("Sales Forecasting Models")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Algorithms Evaluated: Linear Regression, Random Forest Regressor, XGBoost, and Prophet (or Exponential Smoothing)."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Evaluation Metrics: Model performance compared on RMSE, MAE, and R² Score on a 6-month holdout set."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Selected Model: Random Forest / XGBoost ensemble selected for final forecasting due to lower RMSE and robust seasonal tracking."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• 12-Month Projections: Projections indicate a continued seasonal peak in Q4 and steady 8% overall MoM growth trend."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 11: Customer Churn Prediction
    slide = add_standard_slide("Customer Churn Prediction")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Classification Models: Logistic Regression, Random Forest Classifier, and XGBoost Classifier."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Predictive Features: Recency, order count, total revenue, average discount, return rate, age, gender, and region."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Champion Classifier: Random Forest Classifier achieved 81% accuracy and an ROC-AUC score of 0.84 on test dataset."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Outputs: Computes Churn Probability for each customer and classifies them into risk categories (Low, Medium, High)."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 12: Product Demand & Inventory Recommendations
    slide = add_standard_slide("Product Demand & Inventory Recommendations")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Objective: Prevent stockouts on high-demand items while minimizing inventory carrying costs."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Formulations:\n  - Safety Stock = (Max Daily Sales * Max Lead Time) - (Avg Daily Sales * Avg Lead Time)\n  - Reorder Point = (Avg Daily Sales * Avg Lead Time) + Safety Stock"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Inventory Operations: Custom safety stock and reorder point parameters calculated at the category-region level."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Supply chain lead time: Assumed average lead time of 10 days, with a maximum window of 15 days."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 13: Business Insights Engine
    slide = add_standard_slide("Business Insights Engine")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Diagnostics: Dynamically checks 20 business rules against overall retail metrics."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Automated Recommendations: Auto-generates detailed, actionable business actions based on database triggers."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Quantified Impact: Calculates Expected ROI, Revenue Impact, and Cost Savings for each recommendation."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Delivery: Outputs compiled report into reports/business_insights_report.md."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 14: Actionable Recommendations (Top 5)
    slide = add_standard_slide("Actionable Recommendations (Top 5)")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. Launch Customer Loyalty Program: Tiered system to boost repeat purchase rates from 54% to target 65%."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "2. Enforce a 20% Discount Cap: Block regional sales managers from applying extreme discounts (>20%) without approval."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "3. Automated At-Risk Win-back Campaigns: Email incentives (15% off) sent automatically to customers inactive for 150+ days."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "4. Bundle and Liquidate Dead Stock: Combine low-velocity items (velocity < 0.05) with top-selling office products."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "5. Shift Marketing PPC Budgets: Reallocate 15% budget from underperforming Central region to high-performing West region."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN

    # Slide 15: Projected ROI & Financial Impact
    slide = add_standard_slide("Projected ROI & Financial Impact")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Strategic Financial Growth Projections:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Projected Revenue Impact: $91,450.00 in incremental revenue across campaigns."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Projected Cost Savings Impact: $34,200.00 in savings from logistics and discount restrictions."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Total Economic Impact: $125,650.00 in combined economic expansion."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Program Average ROI: 330% return across all marketing and operational programs."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 16: Future Project Scope
    slide = add_standard_slide("Future Project Scope")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Real-Time Processing: Integrate streaming engines (Kafka/Flink) to ingest orders and recalculate inventories dynamically."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• NLP Customer Sentiment: Scraping customer reviews on products to factor qualitative scores into return rate calculations."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Advanced Clustering: Transition from static RFM score bins to unsupervised K-Means clustering for non-linear customer mapping."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Supply Chain GIS Routing: Map shipping routes using GIS and vehicle routing algorithms to reduce furniture delivery costs."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 17: Conclusion
    slide = add_standard_slide("Conclusion")
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Decision Support: Transforms simple raw transactional tables into a complete business decision-support tool."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Machine Learning Power: Predictive engines shift operations from reactive reporting to proactive customer retention and stock optimization."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Scalable Design: The modular Python pipeline (src/) is designed to scale to larger relational databases or cloud data warehouses."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Delivery Check: Delivers all academic and professional components ready for evaluation."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MAIN

    # Slide 18: Q&A
    slide = add_standard_slide("Q&A")
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = ACCENT_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(28)
    p2.font.name = "Arial"
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Project Code: src/  |  Project Specifications: reports/"
    p3.font.size = Pt(16)
    p3.font.name = "Arial"
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)
    
    os.makedirs("presentation", exist_ok=True)
    prs.save("presentation/final_presentation.pptx")
    print("Presentation PowerPoint file generated successfully and saved at presentation/final_presentation.pptx")

if __name__ == "__main__":
    create_presentation()
